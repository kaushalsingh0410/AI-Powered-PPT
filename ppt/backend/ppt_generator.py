"""
PowerPoint Presentation Generator from Python List Data
Converts structured slide data into professional PPTX presentations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTGenerator:
    def __init__(self, output_filename="presentation.pptx",title='presentation.pptx'):
        """Initialize the presentation generator"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.output_filename = output_filename
        self.title = title
        
    def add_title_slide(self, title, subtitle=""):
        """Add a title slide to the presentation"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        return slide
    
    def add_content_slide(self, title, content_items):
        """
        Add a content slide with title and bullet points
        
        Args:
            title (str): Slide title
            content_items (list): List of content items to display as bullets
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content as bullet points
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()  # Clear default text
        
        for idx, item in enumerate(content_items):
            if idx == 0:
                # First paragraph
                p = text_frame.paragraphs[0]
            else:
                # Add new paragraphs
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = 0  # Main level (no indentation)
            p.font.size = Pt(18)
            p.space_after = Pt(12)
        
        return slide
    
    def add_detailed_content_slide(self, title, detailed_content_list):
        """
        Add a slide with detailed content including key points and explanations
        
        Args:
            title (str): Slide title
            detailed_content_list (list): List of dicts with keys like 'key_point', 'explanation', etc.
        """
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for idx, content_item in enumerate(detailed_content_list):
            # Add key point as main bullet
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            key_point = content_item.get("key_point", "")
            p.text = key_point
            p.level = 0
            p.font.size = Pt(16)
            p.font.bold = True
            p.space_after = Pt(6)
            
            # Add additional details as sub-bullets if they exist
            if "explanation" in content_item:
                p_exp = text_frame.add_paragraph()
                p_exp.text = content_item["explanation"]
                p_exp.level = 1
                p_exp.font.size = Pt(14)
                p_exp.space_after = Pt(6)
            
            if "example" in content_item:
                p_ex = text_frame.add_paragraph()
                p_ex.text = f"Example: {content_item['example']}"
                p_ex.level = 1
                p_ex.font.size = Pt(14)
                p_ex.space_after = Pt(6)
            
            if "statistic" in content_item:
                p_stat = text_frame.add_paragraph()
                p_stat.text = f"📊 {content_item['statistic']}"
                p_stat.level = 1
                p_stat.font.size = Pt(14)
                p_stat.space_after = Pt(10)
            
            if "additional_info" in content_item:
                p_info = text_frame.add_paragraph()
                p_info.text = content_item["additional_info"]
                p_info.level = 1
                p_info.font.size = Pt(14)
                p_info.space_after = Pt(10)
        
        return slide
    
    def generate_from_list(self, slides_data):
        """
        Generate presentation from a list of slide dictionaries
        
        Args:
            slides_data (list): List of slide dictionaries with:
                - slide_number (int)
                - slide_title (str)
                - detailed_content (list of dicts)
        """
        # Add title slide
        self.add_title_slide(self.title, "Generated By Ritey")
        
        # Add content slides
        for slide_data in slides_data:
            title = slide_data.get("slide_title", "Slide")
            detailed_content = slide_data.get("detailed_content", [])
            
            if detailed_content:
                self.add_detailed_content_slide(title, detailed_content)
        
        # Add conclusion slide
        self.add_title_slide("Thank You!", "Questions?")
    
    def save(self):
        """Save the presentation to file"""
        self.prs.save(self.output_filename)
        print(f"✅ Presentation saved as: {self.output_filename}")
        return self.output_filename


# # Example usage with your data
# if __name__ == "__main__":
    
#     # Your slide data
#     slides_data = [
#         {
#             "slide_number": 1,
#             "slide_title": "Introduction to React",
#             "detailed_content": [
#                 {
#                     "key_point": "React is a powerful JavaScript library designed specifically for building user interfaces, allowing developers to create interactive web applications efficiently.",
#                     "explanation": "Created by Facebook (now Meta) in 2013, React revolutionized web development by introducing a component-based architecture that promotes reusability and scalability. It is widely adopted for dynamic web apps because it enables developers to update parts of a page without reloading the entire site.",
#                     "example": "For instance, companies like Netflix, Airbnb, and Walmart use React to manage complex user interfaces with hundreds of interactive components.",
#                     "statistic": "According to Stack Overflow's 2023 Developer Survey, React is among the top 3 most-used web frameworks globally, with over 40% of developers reporting its use in their projects."
#                 }
#             ]
#         },
#         {
#             "slide_number": 2,
#             "slide_title": "Core Concepts and Features",
#             "detailed_content": [
#                 {
#                     "key_point": "Component-based architecture allows developers to break down applications into reusable, modular pieces, improving code maintainability and scalability.",
#                     "additional_info": "For example, a login form component can be reused across multiple pages without rewriting the code. This modularity is widely adopted in frameworks like React and Vue, where components encapsulate both logic and UI."
#                 },
#                 {
#                     "key_point": "The Virtual DOM enhances performance by creating a lightweight copy of the actual DOM.",
#                     "additional_info": "When changes occur, it calculates the minimal updates needed and applies them efficiently, reducing lag and resource usage. Studies show this can improve rendering speed by up to 50% compared to direct DOM manipulation."
#                 },
#                 {
#                     "key_point": "JSX syntax bridges the gap between HTML and JavaScript by allowing developers to write HTML-like code directly within JavaScript.",
#                     "additional_info": "For instance, a button element can be defined as <button>Submit</button> without switching languages. JSX is processed by transpilers like Babel into standard JavaScript."
#                 }
#             ]
#         },
#         {
#             "slide_number": 3,
#             "slide_title": "Use Cases and Ecosystem",
#             "detailed_content": [
#                 {
#                     "key_point": "Enterprise applications favor this technology due to its scalability and maintainability.",
#                     "additional_info": "It makes it ideal for large-scale projects where long-term sustainability and adaptability are critical. Companies like Airbnb and Netflix have leveraged similar ecosystems to manage complex user interfaces."
#                 },
#                 {
#                     "key_point": "Integration with popular front-end tools like Redux, React Router, and Next.js enhances developer productivity.",
#                     "additional_info": "Redux alone is used in over 60% of React projects according to recent surveys. These tools handle state management, navigation, and server-side rendering."
#                 },
#                 {
#                     "key_point": "A strong community and extensive documentation solidify its value.",
#                     "additional_info": "Developers have access to tutorials, forums, and open-source libraries that accelerate onboarding. Over 1.2 million repositories on GitHub reference React technology."
#                 }
#             ]
#         }
#     ]
    
#     # Create and generate presentation
#     generator = PPTGenerator("React_Presentation.pptx")
#     generator.generate_from_list(slides_data)
#     generator.save()
