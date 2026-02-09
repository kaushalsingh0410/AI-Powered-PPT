"""
Beautiful PowerPoint Presentation Generator from Python List Data
Creates professionally styled presentations with modern design, colors, and layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class BeautifulPPTGenerator:
    def __init__(self, output_filename="presentation.pptx"):
        """Initialize the presentation generator with beautiful styling"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.output_filename = output_filename
        
        # Color Palette - Modern & Professional
        self.PRIMARY_COLOR = RGBColor(26, 35, 126)      # Deep Blue
        self.ACCENT_COLOR = RGBColor(76, 175, 255)      # Sky Blue
        self.SECONDARY_COLOR = RGBColor(244, 67, 54)    # Modern Red
        self.DARK_BG = RGBColor(245, 247, 250)          # Light Gray
        self.TEXT_PRIMARY = RGBColor(33, 33, 33)        # Dark Gray
        self.TEXT_SECONDARY = RGBColor(117, 117, 117)   # Medium Gray
        self.WHITE = RGBColor(255, 255, 255)            # White
        
    def _set_background(self, slide, color):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
    
    def _add_decorative_line(self, slide, top_inches, color):
        """Add a decorative line to the slide"""
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(0.5), Inches(top_inches),
            Inches(9), Inches(0)
        )
        line.line.color.rgb = color
        line.line.width = Pt(3)
    
    def add_beautiful_title_slide(self, title, subtitle=""):
        """Add a stunning title slide with modern design"""
        # Create blank slide
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Set background
        self._set_background(slide, self.PRIMARY_COLOR)
        
        # Add decorative shape
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT_COLOR
        shape.line.color.rgb = self.ACCENT_COLOR
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(9), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.8),
            Inches(9), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        
        p_sub = subtitle_frame.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(28)
        p_sub.font.color.rgb = self.WHITE
        p_sub.alignment = PP_ALIGN.CENTER
        
        # Add decorative line at bottom
        self._add_decorative_line(slide, 6.8, self.SECONDARY_COLOR)
        
        return slide
    
    def add_beautiful_content_slide(self, title, content_items):
        """Add a content slide with beautiful formatting"""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Set background
        self._set_background(slide, self.WHITE)
        
        # Add header bar
        header = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_COLOR
        header.line.color.rgb = self.PRIMARY_COLOR
        
        # Add decorative accent
        accent_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0.95),
            Inches(10), Inches(0.08)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        accent_bar.line.color.rgb = self.ACCENT_COLOR
        
        # Add title in header
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.2),
            Inches(8.6), Inches(0.7)
        )
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_title = title_frame.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(40)
        p_title.font.bold = True
        p_title.font.color.rgb = self.WHITE
        
        # Add content area
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.5),
            Inches(8), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for idx, item in enumerate(content_items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = "◆ " + item  # Add beautiful bullet
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.TEXT_PRIMARY
            p.space_before = Pt(8)
            p.space_after = Pt(12)
        
        return slide
    
    def add_beautiful_detailed_slide(self, title, detailed_content_list):
        """Add a detailed content slide with beautiful formatting"""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Set background
        self._set_background(slide, self.DARK_BG)
        
        # Add header bar
        header = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(0.9)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.PRIMARY_COLOR
        header.line.color.rgb = self.PRIMARY_COLOR
        
        # Add decorative accent
        accent_bar = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0.85),
            Inches(10), Inches(0.08)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        accent_bar.line.color.rgb = self.ACCENT_COLOR
        
        # Add title in header
        title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.15),
            Inches(8.6), Inches(0.65)
        )
        title_frame = title_box.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_title = title_frame.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(38)
        p_title.font.bold = True
        p_title.font.color.rgb = self.WHITE
        
        # Add content with cards
        current_top = 1.3
        card_height = 1.2
        
        for idx, content_item in enumerate(detailed_content_list):
            if current_top + card_height > 7.0:
                break
            
            # Add card background
            card = slide.shapes.add_shape(
                1,
                Inches(0.7), Inches(current_top),
                Inches(8.6), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.WHITE
            card.line.color.rgb = self.ACCENT_COLOR
            card.line.width = Pt(2)
            
            # Add key point (bold)
            key_point_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(current_top + 0.1),
                Inches(8.2), Inches(0.4)
            )
            kp_frame = key_point_box.text_frame
            kp_frame.word_wrap = True
            
            p_kp = kp_frame.paragraphs[0]
            p_kp.text = "★ " + content_item.get("key_point", "")[:100] + "..."
            p_kp.font.size = Pt(14)
            p_kp.font.bold = True
            p_kp.font.color.rgb = self.PRIMARY_COLOR
            
            # Add additional info
            info_text = ""
            if "explanation" in content_item:
                info_text = content_item["explanation"][:80] + "..."
            elif "additional_info" in content_item:
                info_text = content_item["additional_info"][:80] + "..."
            
            if info_text:
                info_box = slide.shapes.add_textbox(
                    Inches(0.9), Inches(current_top + 0.55),
                    Inches(8.2), Inches(0.6)
                )
                info_frame = info_box.text_frame
                info_frame.word_wrap = True
                
                p_info = info_frame.paragraphs[0]
                p_info.text = info_text
                p_info.font.size = Pt(12)
                p_info.font.color.rgb = self.TEXT_SECONDARY
            
            current_top += card_height + 0.15
        
        return slide
    
    def add_beautiful_closing_slide(self):
        """Add a beautiful closing/thank you slide"""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Set background with gradient effect using shapes
        self._set_background(slide, self.PRIMARY_COLOR)
        
        # Add accent shape
        accent_shape = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(10), Inches(3)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = self.ACCENT_COLOR
        accent_shape.line.color.rgb = self.ACCENT_COLOR
        
        # Add thank you text
        thank_you_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(8), Inches(1.5)
        )
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_thank = thank_you_frame.paragraphs[0]
        p_thank.text = "Thank You!"
        p_thank.font.size = Pt(60)
        p_thank.font.bold = True
        p_thank.font.color.rgb = self.WHITE
        p_thank.alignment = PP_ALIGN.CENTER
        
        # Add questions text
        questions_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(8), Inches(1)
        )
        questions_frame = questions_box.text_frame
        questions_frame.word_wrap = True
        
        p_q = questions_frame.paragraphs[0]
        p_q.text = "Questions & Discussion"
        p_q.font.size = Pt(32)
        p_q.font.color.rgb = self.PRIMARY_COLOR
        p_q.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def generate_from_list(self, slides_data):
        """Generate beautiful presentation from a list of slide dictionaries"""
        # Add stunning title slide
        self.add_beautiful_title_slide(
            "Presentation",
            "✨ Generated from Python Data"
        )
        
        # Add content slides
        for slide_data in slides_data:
            title = slide_data.get("slide_title", "Slide")
            detailed_content = slide_data.get("detailed_content", [])
            
            if detailed_content:
                self.add_beautiful_detailed_slide(title, detailed_content)
        
        # Add beautiful closing slide
        self.add_beautiful_closing_slide()
    
    def save(self):
        """Save the presentation to file"""
        self.prs.save(self.output_filename)
        print(f"✨ Beautiful presentation saved as: {self.output_filename}")
        return self.output_filename


# Example usage with your data
if __name__ == "__main__":
    
    # Your slide data
    slides_data = [
        {
            "slide_number": 1,
            "slide_title": "Introduction to React",
            "detailed_content": [
                {
                    "key_point": "React is a powerful JavaScript library designed specifically for building user interfaces, allowing developers to create interactive web applications efficiently.",
                    "explanation": "Created by Facebook (now Meta) in 2013, React revolutionized web development by introducing a component-based architecture that promotes reusability and scalability. It is widely adopted for dynamic web apps because it enables developers to update parts of a page without reloading the entire site.",
                    "example": "For instance, companies like Netflix, Airbnb, and Walmart use React to manage complex user interfaces with hundreds of interactive components.",
                    "statistic": "According to Stack Overflow's 2023 Developer Survey, React is among the top 3 most-used web frameworks globally, with over 40% of developers reporting its use in their projects."
                }
            ]
        },
        {
            "slide_number": 2,
            "slide_title": "Core Concepts and Features",
            "detailed_content": [
                {
                    "key_point": "Component-based architecture allows developers to break down applications into reusable, modular pieces, improving code maintainability and scalability.",
                    "additional_info": "For example, a login form component can be reused across multiple pages without rewriting the code. This modularity is widely adopted in frameworks like React and Vue, where components encapsulate both logic and UI."
                },
                {
                    "key_point": "The Virtual DOM enhances performance by creating a lightweight copy of the actual DOM.",
                    "additional_info": "When changes occur, it calculates the minimal updates needed and applies them efficiently, reducing lag and resource usage. Studies show this can improve rendering speed by up to 50% compared to direct DOM manipulation."
                },
                {
                    "key_point": "JSX syntax bridges the gap between HTML and JavaScript by allowing developers to write HTML-like code directly within JavaScript.",
                    "additional_info": "For instance, a button element can be defined as <button>Submit</button> without switching languages. JSX is processed by transpilers like Babel into standard JavaScript."
                }
            ]
        },
        {
            "slide_number": 3,
            "slide_title": "Use Cases and Ecosystem",
            "detailed_content": [
                {
                    "key_point": "Enterprise applications favor this technology due to its scalability and maintainability.",
                    "additional_info": "It makes it ideal for large-scale projects where long-term sustainability and adaptability are critical. Companies like Airbnb and Netflix have leveraged similar ecosystems to manage complex user interfaces."
                },
                {
                    "key_point": "Integration with popular front-end tools like Redux, React Router, and Next.js enhances developer productivity.",
                    "additional_info": "Redux alone is used in over 60% of React projects according to recent surveys. These tools handle state management, navigation, and server-side rendering."
                },
                {
                    "key_point": "A strong community and extensive documentation solidify its value.",
                    "additional_info": "Developers have access to tutorials, forums, and open-source libraries that accelerate onboarding. Over 1.2 million repositories on GitHub reference React technology."
                }
            ]
        }
    ]
    
    # Create and generate beautiful presentation
    generator = BeautifulPPTGenerator("React_Presentation_Beautiful.pptx")
    generator.generate_from_list(slides_data)
    generator.save()
